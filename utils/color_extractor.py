from PIL import Image
import numpy as np
from sklearn.cluster import KMeans
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def extract_colors(image_path, n_colors=3):
    img = Image.open(image_path).convert("RGB")
    img = img.resize((150, 150))
    arr = np.array(img).reshape((-1, 3))
    
    k_means = KMeans(n_clusters=n_colors, random_state=42)
    k_means.fit(arr)
    colors = k_means.cluster_centers_.astype(int)
    
    return [tuple(c) for c in colors]

def get_contrast_color(rgb):
    r, g, b = rgb
    luminance = (0.299*r + 0.587*g + 0.114*b)/255
    return (0, 0, 0) if luminance > 0.5 else (255, 255, 255)

def create_theme_ppt(colors, output_path="brand_theme.pptx"):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    bg_color = colors[0]
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*map(int, bg_color))
    
    title = slide.shapes.title
    title.text = "Brand Title Slide"
    
    # Set title font color to contrast color
    contrast_rgb = get_contrast_color(bg_color)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*contrast_rgb)
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content slide
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    
    bg_color2 = colors[1]
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = RGBColor(*map(int, bg_color2))
    
    title2 = slide2.shapes.title
    title2.text = "Content Slide"
    contrast_rgb2 = get_contrast_color(bg_color2)
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(*contrast_rgb2)
    
    content = slide2.placeholders[1]
    content.text = "• Point 1\n• Point 2\n• Point 3"
    for para in content.text_frame.paragraphs:
        para.font.size = Pt(24)
        para.font.color.rgb = RGBColor(*contrast_rgb2)
    
    prs.save(output_path)
    print(f"PowerPoint theme saved to {output_path}")

if __name__ == "__main__":
    logo_path = r"C:\Users\ACER\Documents\GitHub\bird-colorful-logo-gradient-vector_343694-1365.jpg"
    colors = extract_colors(logo_path, n_colors=3)
    print("Extracted colors:", colors)
    
    create_theme_ppt(colors)
